# diffcult_finish.py – AI Teaching Assistant with Authentication & Difficulty Levels
# -----------------------------------------------------------------------------
# Combines the authentication / user‑management features from **app.py** with the
# difficulty‑aware question workflow from **diffcult_finish.py** and now shows
# full multiple‑choice options right beside each radio button.
# -----------------------------------------------------------------------------
# ✨ 2025‑07‑14  – Update:
#    • Added extract_mc_options() to parse and display full option text.
#    • Radio widgets for Q6‑Q10 now list "A) option", "B) option", … instead of
#      just the letters.
#    • Grading logic updated to map the selected label back to its leading letter.
# -----------------------------------------------------------------------------

import streamlit as st
import google.generativeai as genai
import fitz, os, re, json, random, hashlib
from datetime import datetime, timedelta
from docx import Document
from pptx import Presentation
#import config  # GEMINI_API_KEY & GEMINI_MODEL_ID must be defined here
import os
genai.configure(api_key=os.environ["GEMINI_API_KEY"])
MODEL_ID = os.environ["gemini-pro"]

# ─────────────────────────────── CONSTANTS ────────────────────────────────
USERS_FILE    = "users.json"
WRONG_LOG     = "wrong_log.txt"
MATERIAL_FILE = "material.txt"
LOCKOUT_MIN   = 5    # minutes (after 3 wrong PIN attempts)
MAX_PIN_TRY   = 3

# ─────────────────────────── Gemini Settings ──────────────────────────────
genai.configure(api_key=config.GEMINI_API_KEY)
MODEL_ID = config.GEMINI_MODEL_ID

# ╔═══════════════════ USER / AUTHENTICATION UTILITIES ═══════════════════╗

def load_users() -> dict:
    if not os.path.exists(USERS_FILE):
        return {}
    try:
        with open(USERS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

def save_users(data: dict):
    with open(USERS_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)

def hash_pw(pwd: str) -> str:
    return hashlib.sha256(pwd.encode()).hexdigest()

# ╚═══════════════════════════════════════════════════════════════════════╝

# ╔════════════════════ FILE / MATERIAL UTILITIES ════════════════════════╗

def load_material() -> str:
    if os.path.exists(MATERIAL_FILE):
        with open(MATERIAL_FILE, "r", encoding="utf-8") as f:
            return f.read()
    return ""


def extract_text(uploaded_file):
    """Return plain text from PDF / DOCX / PPTX."""
    ext = uploaded_file.name.split(".")[-1].lower()
    if ext == "pdf":
        return "\n".join(p.get_text() for p in fitz.open(stream=uploaded_file.read(), filetype="pdf"))
    if ext == "docx":
        return "\n".join(p.text for p in Document(uploaded_file).paragraphs)
    if ext == "pptx":
        return "\n".join(
            shp.text for slide in Presentation(uploaded_file).slides for shp in slide.shapes if getattr(shp, "text", "")
        )
    st.warning("Unsupported file type"); return ""

# ╚═══════════════════════════════════════════════════════════════════════╝

# ╔════════════════════ WRONG‑LOG UTILITIES ══════════════════════════════╗

def clean_json(raw: str) -> dict:
    """Attempt to parse *anything* that looks like JSON from LLM output."""
    try:
        return json.loads(raw)
    except Exception:
        pass
    m = re.search(r"\{.*\}", raw, re.S)
    if m:
        try:
            return json.loads(m.group())
        except Exception:
            pass
    return {}


def append_wrong_log(username: str, level: str, wrong: list, exps: dict):
    """Append wrong answers to WRONG_LOG with per‑user & difficulty tags."""
    lines = []
    for w in wrong:
        lines += [
            f"{username} | {w['id']} [{level}]: {w['question']}",
            f"Your answer: {w['got']} | Correct: {w['expected']}",
            f"Explanation: {exps.get(w['id'], '(no explanation)')}",
            "----",
        ]
    with open(WRONG_LOG, "a", encoding="utf-8") as f:
        f.write("\n".join(lines) + "\n")


def load_wrong_log(username: str = None):
    """Return list of wrong‑question dicts (optionally filtered by username)."""
    if not os.path.exists(WRONG_LOG):
        return []
    entries = []
    with open(WRONG_LOG, "r", encoding="utf-8") as f:
        for block in filter(None, f.read().split("----")):
            lines = [l.strip() for l in block.strip().splitlines() if l.strip()]
            if len(lines) < 3:
                continue
            # Username | Q7 [Hard]: ………………………
            m = re.match(r"([^|]+)\|\s*(Q\d+)\s*\[\s*(Easy|Medium|Hard)\s*\]:\s*(.*)", lines[0])
            if not m:
                continue
            uname, qid, diff, qtxt = m.groups()
            if username and uname.strip() != username:
                continue
            m_cor = re.search(r"Correct:\s*([A-Za-z]+)", lines[1])
            correct = m_cor.group(1).upper() if m_cor else "UNKNOWN"
            entries.append({
                "id":          qid,
                "question":    qtxt,
                "correct":     correct,
                "explanation": lines[2].replace("Explanation:", "").strip(),
                "difficulty":  diff.capitalize(),
            })
    return entries

# ╚═══════════════════════════════════════════════════════════════════════╝

# ╔════════════════════ MISC. STREAMLIT HELPERS ══════════════════════════╗

def clear_radio_keys():
    for k in list(st.session_state):
        if k.startswith(("tf", "mc")):
            del st.session_state[k]

def do_rerun():
    if hasattr(st, "rerun"):
        st.rerun()
    else:
        st.experimental_rerun()


def extract_mc_options(question: str):
    """Extract multiple‑choice texts and return a list like ['A) …', 'B) …', …]."""
    # Normalize full‑width parentheses
    q = question.replace("（", "(").replace("）", ")")
    matches = re.findall(r"\([A-D]\)\s*([^()]+)", q)
    if len(matches) >= 4:
        return [f"{chr(65+i)}) {opt.strip()}" for i, opt in enumerate(matches[:4])]
    return ["A", "B", "C", "D"]  # Fallback if parsing fails

# ╚═══════════════════════════════════════════════════════════════════════╝

# ─────────────────────────── STREAMLIT INIT ──────────────────────────────
st.set_page_config(page_title="AI Teaching Assistant")
if "material" not in st.session_state:
    st.session_state["material"] = load_material()
users = load_users()

# ╔═══════════════════════ AUTHENTICATION UI ═════════════════════════════╗
if "user" not in st.session_state:
    st.title("🔐 Login · Register · Forgot Password")
    act = st.radio("Select option", ("Login", "Register", "Forgot Password"))

    if act == "Login":
        u = st.text_input("Username")
        p = st.text_input("Password", type="password")
        if st.button("Login"):
            if u in users and users[u]["password"] == hash_pw(p):
                st.session_state.update({"user": u, "role": users[u]["role"]}); do_rerun()
            else:
                st.error("Invalid username or password.")

    elif act == "Register":
        role = st.selectbox("Register as", ("Teacher", "Student"))
        u = st.text_input("Username")
        p = st.text_input("Password", type="password")
        pin = st.text_input("4‑digit PIN", max_chars=4)
        if st.button("Create account"):
            if not (u and p and pin):
                st.error("All fields required.")
            elif u in users:
                st.error("Username already exists.")
            elif not pin.isdigit() or len(pin) != 4:
                st.error("PIN must be 4 digits.")
            else:
                users[u] = {"password": hash_pw(p), "role": role, "pin": pin, "fail": 0, "lock": ""}
                save_users(users)
                st.success("Account created – please log in.")

    else:  # Forgot Password / PIN recovery
        u = st.text_input("Username")
        pin = st.text_input("4‑digit PIN", max_chars=4)
        if st.button("Recover"):
            if u not in users:
                st.error("Username not found.")
            else:
                data = users[u]
                locked = data.get("lock", "")
                if locked and datetime.utcnow() < datetime.fromisoformat(locked):
                    wait = datetime.fromisoformat(locked) - datetime.utcnow()
                    st.warning(f"Locked. Try again in {wait.seconds//60}m {wait.seconds%60}s")
                elif pin == data.get("pin"):
                    data.update({"fail": 0, "lock": ""}); save_users(users)
                    st.info("Password SHA‑256 hash (reset via users.json if needed):")
                    st.code(data["password"], language="text")
                else:
                    data["fail"] = data.get("fail", 0) + 1
                    if data["fail"] >= MAX_PIN_TRY:
                        data.update({"lock": (datetime.utcnow() + timedelta(minutes=LOCKOUT_MIN)).isoformat(), "fail": 0})
                        st.error(f"Wrong PIN ×{MAX_PIN_TRY}. Locked {LOCKOUT_MIN} minutes.")
                    else:
                        st.error("Incorrect PIN."); save_users(users)
    st.stop()
# ╚═══════════════════════════════════════════════════════════════════════╝

# ───────────────────────────── MAIN APP ──────────────────────────────────
username, role = st.session_state["user"], st.session_state["role"]
st.sidebar.write(f"👤 **{username}** ({role})")
if st.sidebar.button("Log out"):
    st.session_state.clear(); do_rerun()

st.title("AI Teaching Assistant – Circuit Theory")

# ╔══════════════════════════ TEACHER VIEW ═══════════════════════════════╗
if role == "Teacher":
    upl = st.file_uploader("Upload material (PDF / DOCX / PPTX)", type=["pdf", "docx", "pptx"])
    if upl:
        txt = extract_text(upl)
        st.text_area("Preview extracted text", txt, height=250)
        if st.button("✅ Save material for students"):
            with open(MATERIAL_FILE, "w", encoding="utf-8") as f:
                f.write(txt)
            st.session_state["material"] = txt
            st.success("Material saved & shared with students.")
# ╚═══════════════════════════════════════════════════════════════════════╝

# ╔══════════════════════════ STUDENT VIEW ═══════════════════════════════╗
else:
    if not st.session_state.get("material"):
        st.warning("No material uploaded yet. Please ask your teacher."); st.stop()

    func = st.radio("Select mode", ("📄 New question practice", "🔄 Wrong‑question practice"))

    # ─── 1. New‑question practice ─────────────────────────────────────────
    if func.startswith("📄"):
        material = st.session_state["material"]
        difficulty = st.selectbox("Choose question difficulty", ("Easy", "Medium", "Hard"), index=1)

        if st.button("Generate 10 questions"):
            clear_radio_keys()
            #debug_exp = st.expander("LLM raw output (debug)", expanded=False)

            diff_word = {"Easy": "easy", "Medium": "intermediate", "Hard": "challenging"}[difficulty]

            # Retry loop in case LLM returns malformed output
            for _ in range(5):
                q_prompt = (
                    f"Generate **exactly 10 lines** of {diff_word}-level questions only, following this format:\n"
                    "• Lines 1‑5 start with **T/F** (True/False statements).\n"
                    "• Lines 6‑10 start with **Q6.** … **Q10.** each with options (A)…(D)…\n"
                    "Do NOT include answers, explanations, or duplicate sets.\n\n"
                    f"Material:\n{material}"
                )
                raw_q = genai.GenerativeModel(MODEL_ID).generate_content(q_prompt).text
                #debug_exp.write(raw_q)

                q_lines = []
                for raw_ln in raw_q.splitlines():
                    ln = re.sub(r"^[\s\u2022\*\-]+", "", raw_ln).strip()
                    ln = re.sub(r"^T/F\s*:?", "T/F", ln, flags=re.I)
                    ln = re.sub(r"^Q(\d+)\s*\.?", r"Q\1.", ln, flags=re.I)
                    if re.match(r"^T/F", ln, re.I) or re.match(r"^Q\d+\.", ln, re.I):
                        q_lines.append(ln)
                    if len(q_lines) == 10:
                        break

                valid = (
                    len(q_lines) == 10 and
                    all(ln.upper().startswith("T/F") for ln in q_lines[:5]) and
                    all(re.match(r"^Q(?:[6-9]\.|10\.)", ln, re.I) for ln in q_lines[5:])
                )
                if not valid:
                    continue

                # Generate answer key
                ak_prompt = (
                    "Provide ONLY a JSON mapping {'Q1': 'True'|'False', …, 'Q10': 'A'..'D'}.\n"
                    "Ensure at least one of Q6‑Q10 is 'A'.\n\nQuestions:\n" + "\n".join(q_lines)
                )
                answer_key = clean_json(genai.GenerativeModel(MODEL_ID).generate_content(ak_prompt).text)

                if len(answer_key) == 10 and any(
                        answer_key.get(f"Q{i}", "").upper() == "A" for i in range(6, 11)):
                    break
            else:
                st.error("Generation failed. Check the debug panel."); st.stop()

            display_lines = []
            for ln in q_lines:
                # 題幹去掉 (A…)(B…)(C…)(D…)
                if ln.upper().startswith("Q"):
                    ln = re.split(r"\s*\(A", ln, maxsplit=1)[0].rstrip()
                display_lines.append(ln)

            st.session_state.update({
                "gen_id": st.session_state.get("gen_id", 0) + 1,
                "questions": q_lines,          # ← 保留完整文本（含選項）
                "display_q": display_lines,    # ← 顯示用題幹
                "answer_key": answer_key,
                "mode": "new",
                "difficulty": difficulty,
            })
            st.success("Questions generated!")
        gid = st.session_state.get("gen_id", 0)
        
        if st.session_state.get("mode") == "new":
            stems = st.session_state.get("display_q", st.session_state["questions"])

            for idx, (raw_q, stem) in enumerate(
                    zip(st.session_state["questions"], stems), 1):

                # 只顯示題幹（無選項）
                st.markdown(f"**{stem}**")

                # 依題型產生 Radio
                qid = f"Q{idx}"
                if raw_q.upper().startswith("T/F"):
                    st.radio("", ["True", "False"], key=f"tf_{qid}_{gid}")
                else:
                    opts = re.findall(r"\((A|B|C|D)\)\s*([^()]+)", raw_q)
                    labels = [f"{ltr}) {txt.strip()}" for ltr, txt in opts]
                    st.radio("", labels, key=f"mc_{qid}_{gid}")

    # ─── 2. Wrong‑question practice ───────────────────────────────────────
    else:
        level = st.selectbox("Choose difficulty for review", ("Easy", "Medium", "Hard", "All"), index=3)

        if st.button("Start wrong‑question practice"):
            wbank = load_wrong_log(username=username)
            if level != "All":
                wbank = [e for e in wbank if e.get("difficulty") == level]

            if not wbank:
                st.error(f"No wrong questions recorded for “{level}” level."); st.stop()
            
            if len(wbank) < 10:
                st.error("Not enough wrong questions"); st.stop()
            
            uniq_map = {}
            for e in wbank:
                uniq_map[e["id"]] = e  # 保留最後一次出現的錯題

            sample = random.sample(list(uniq_map.values()), min(10, len(uniq_map)))
            qs   = [(e["id"], e["question"]) for e in sample]
            ak   = {e["id"]: e["correct"] for e in sample}
            expl = {e["id"]: e["explanation"] for e in sample}

            st.session_state.update({
                "gen_id":       st.session_state.get("gen_id", 0) + 1,
                "questions":    qs,
                "answer_key":   ak,
                "explanations": expl,
                "mode":         "wrong",
            })
            st.success("Questions loaded – start answering!")

# ────────────────────────── QUIZ DISPLAY & SCORING ───────────────────────
qs   = st.session_state.get("questions")
ak   = st.session_state.get("answer_key")
gid  = st.session_state.get("gen_id", 0)
mode = st.session_state.get("mode", "")

if qs and ak:
    if mode == "wrong":
        st.header("Answer the questions")
        gid = st.session_state.get("gen_id", 0)
        for qid, qtxt in qs:
            num = int(qid[1:])
            if num <= 5:
                st.markdown(f"**{qtxt}**")
                st.radio("", ["True", "False"], key=f"wrong_tf_{qid}_{gid}")
            else:
                stem = re.split(r"\s*\(A", qtxt, maxsplit=1)[0].rstrip()
                opts = re.findall(r"\((A|B|C|D)\)\s*([^()]+)", qtxt)
                labels = [f"{ltr}) {txt.strip()}" for ltr, txt in opts]
                st.markdown(f"**{stem}**")
                st.radio("", labels, key=f"wrong_mc_{qid}_{gid}")
                
    if st.button("Submit answers"):
        wrong = []
        tf_map = {"T": "TRUE", "F": "FALSE"}

        for idx, qitem in enumerate(qs, 1):
            qid, qtxt = (qitem if isinstance(qitem, tuple) else (f"Q{idx}", qitem))
            num = int(qid[1:])
            prefix = "wrong_" if mode == "wrong" else ""
            sid = f"{prefix}tf_{qid}_{gid}" if num <= 5 else f"{prefix}mc_{qid}_{gid}"
            stud_raw = str(st.session_state.get(sid, ""))

            if num <= 5:
                stud = tf_map.get(stud_raw.upper(), stud_raw.upper())
            else:
                stud = stud_raw[:1].upper() if stud_raw else ""

            corr = ak.get(qid, "").strip().upper()
            corr = tf_map.get(corr, corr)
            if stud != corr:
                wrong.append({"id": qid, "question": qtxt, "expected": corr, "got": stud})

        score = (10 - len(wrong)) * 10
        st.subheader(f"Score: {score} / 100")

        # Generate or retrieve explanations
        expl = {}
        if wrong:
            if mode == "new":
                mat = st.session_state.get("material", "")
                prompt = (
                    "For each question below, give a concise explanation "
                    "(1‑3 sentences) of why the correct answer is right. "
                    "Respond JSON id→explanation.\n\n" +
                    f"Material:\n{mat}\n\n" +
                    "\n".join(f"{w['id']}: {w['question']}" for w in wrong)
                )
                expl = clean_json(genai.GenerativeModel(MODEL_ID, generation_config={"temperature": 0.3}).generate_content(prompt).text)
            else:
                expl = st.session_state.get("explanations", {})

        # Display feedback
        if wrong:
            st.write("### Incorrect answers & explanations")
            for w in wrong:
                st.markdown(
                    f"**{w['id']}** – expected **{w['expected']}**, you chose **{w['got']}**  \n"
                    f"> {expl.get(w['id'], '(no explanation)')}"
                )
        else:
            st.success("Perfect! All answers correct.")

        # Log wrong answers (only in new‑question mode)
        if wrong and mode == "new":
            lvl = st.session_state.get("difficulty", "Medium")
            append_wrong_log(username, lvl, wrong, expl)
            st.info(f"Appended {len(wrong)} question(s) to wrong‑log for user {username}.")
